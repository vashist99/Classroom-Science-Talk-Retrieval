"""Generate the Dr. Hadley results presentation as a .pptx file.

Builds a concise yet comprehensive deck summarizing the science-talk-detection
pipeline: the architecture in plain terms, the routing-audit results, the
caveat about the small/skewed seed corpus, and the Y2-vs-more-labels decision.

Run from the project root:
    python scripts/build_hadley_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PATH = PROJECT_ROOT / "Science_Talk_Detection_Hadley_Deck.pptx"
ARCH_IMAGE = PROJECT_ROOT / "architecture_diagram.png"

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette (matches the architecture diagram's blue/amber theme)
NAVY = RGBColor(0x1A, 0x2B, 0x3C)
BLUE = RGBColor(0x2E, 0x5C, 0x8A)
LIGHT_BLUE = RGBColor(0xE8, 0xF1, 0xFB)
AMBER = RGBColor(0xB8, 0x86, 0x2A)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _set_run(run, *, size, bold=False, color=NAVY, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _title_bar(slide, text, *, subtitle=None):
    """Top accent bar + slide title."""
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_W, Inches(0.28)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    _, tf = _textbox(slide, Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.0))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    _set_run(r, size=30, bold=True, color=NAVY)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        _set_run(r2, size=15, italic=True, color=GREY)


def _bullets(slide, items, *, left=Inches(0.7), top=Inches(1.7),
             width=Inches(12.0), height=Inches(5.2), size=18, gap=8):
    """items: list of (text, level, bold, color). Returns the text frame."""
    _, tf = _textbox(slide, left, top, width, height)
    first = True
    for text, level, bold, color in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        bullet = "" if text == "" else ("•  " if level == 0 else "·  ")
        r = p.add_run()
        r.text = f"{bullet}{text}"
        _set_run(r, size=size - (2 if level else 0), bold=bold,
                 color=color or (NAVY if level == 0 else GREY))
    return tf


def _table(slide, rows, *, left, top, width, height, header=True,
           col_widths=None, font=14):
    n_rows = len(rows)
    n_cols = len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(val)
            is_header = header and r_idx == 0
            _set_run(
                run,
                size=font,
                bold=is_header,
                color=WHITE if is_header else NAVY,
            )
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.fore_color.rgb = LIGHT_BLUE if r_idx % 2 else WHITE
    return table


def _note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = _blank(prs)
    _bg(s, NAVY)
    band = s.shapes.add_shape(1, Inches(0), Inches(2.6), SLIDE_W, Inches(0.08))
    band.fill.solid(); band.fill.fore_color.rgb = AMBER
    band.line.fill.background()

    _, tf = _textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.2))
    r = tf.paragraphs[0].add_run()
    r.text = "Detecting Science Talk in Parent and Child Conversations"
    _set_run(r, size=38, bold=True, color=WHITE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Building the training-data engine for an automatic detector"
    _set_run(r, size=20, italic=True, color=RGBColor(0xBF, 0xD4, 0xEA))

    _, tf2 = _textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0))
    r = tf2.paragraphs[0].add_run()
    r.text = "Pipeline Results & Path Forward  ·  Prepared for Dr. Hadley"
    _set_run(r, size=16, color=RGBColor(0x9F, 0xB8, 0xD0))
    _note(s,
        "Opening: \"Thanks for taking the time. In the next ten minutes I want to do "
        "three things: show you how the detector works in plain language, walk through "
        "what your audit and our evaluation actually found, and then ask you for one "
        "decision about where we point this next.\" Keep it light, no jargon yet. Set the "
        "expectation that there is a clear ask at the end so the audience listens for it.")


def slide_problem(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "The Problem")
    _bullets(s, [
        ("Hours of recorded talk, and the science moments are buried in small talk and logistics.", 0, False, None),
        ("Hand-coding every transcript by ear is slow and expensive.", 0, False, None),
        ("Goal: automatically surface and tag the science-talk moments for researchers to review.", 0, True, BLUE),
        ("", 0, False, None),
        ("We are not replacing coders. We hand them a pre-filtered, ranked shortlist.", 1, True, GREY),
    ], top=Inches(2.0), size=22, gap=14)
    _note(s,
        "The pain point first: a single visit produces hours of audio, and the vast "
        "majority of it is logistics, redirection, and small talk. The science moments "
        "are real but sparse, and finding them by ear is the bottleneck that limits how "
        "much data the team can actually use. Emphasize the framing on the last line: this "
        "is an assistant that hands coders a ranked shortlist, and they stay in control of "
        "the final call. We just remove the needle-in-a-haystack search. That framing "
        "matters because it lowers the bar: we do not need perfection, we need a reliable "
        "shortlist.")


def slide_architecture(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "How It Works, in the Simplest Terms")
    if ARCH_IMAGE.exists():
        s.shapes.add_picture(str(ARCH_IMAGE), Inches(1.4), Inches(1.55),
                             width=Inches(10.5))
    _, tf = _textbox(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.6))
    steps = ("1. Record  >  2. Transcribe (speech-to-text)  >  "
             "3. Clean up into sentences  >  4. Detect & label the science sentences")
    r = tf.paragraphs[0].add_run()
    r.text = steps
    _set_run(r, size=18, bold=True, color=BLUE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Each stage hands a clean result to the next, so any piece can be swapped without rewiring the rest."
    _set_run(r, size=14, italic=True, color=GREY)
    _note(s,
        "Walk left to right along the diagram. Stage 1: the parent wears a small recorder, "
        "and nothing is processed on the device. Stage 2: speech-to-text turns the audio "
        "into a written transcript (we use Whisper, chosen because it holds up in noisy "
        "rooms and with child speech). Stage 3: we tidy the text and split it into "
        "individual sentences. Stage 4 is the brain, where the system decides which "
        "sentences are science talk and what kind. The one line to land: each stage hands "
        "a clean result to the next, so we can upgrade any single piece, such as a better "
        "transcriber or a better detector, without rebuilding the rest. Tell them stage 4 "
        "is the next slide.")


def slide_detector(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "The Detector, by Analogy")
    _bullets(s, [
        ("We keep a library of known science-talk examples (\"anchors\").", 0, False, None),
        ("A new sentence is converted into a \"meaning fingerprint.\"", 0, False, None),
        ("We measure how close its fingerprint is to our library.", 0, False, None),
        ("Close enough? Flag it and copy the label of its nearest match.", 0, True, BLUE),
        ("Unsure cases get a second-opinion AI check: \"is this really science talk?\"", 0, True, AMBER),
        ("", 0, False, None),
        ("\"Fingerprint\" = embedding · \"closeness\" = cosine similarity · second opinion = LLM gate.", 1, True, GREY),
    ], top=Inches(1.9), size=21, gap=12)
    _note(s,
        "This is the heart of the talk, so slow down. Analogy: we keep a library of "
        "sentences we know are science talk. When a new sentence comes in, the computer "
        "turns it into a kind of 'meaning fingerprint', a numeric summary of what it "
        "means, not just the words it uses. We then ask: which library example is this "
        "fingerprint closest to? If it is close enough to a known example, we flag it and "
        "borrow that example's label (observation, prediction, and so on). For the "
        "genuinely borderline ones, we get a second opinion from a language model that "
        "reads the sentence and answers, in plain terms, 'is this really science talk?' "
        "Only mention the technical terms on the last line if someone asks: embedding, "
        "cosine similarity, LLM gate. The takeaway: it matches on meaning, not keywords, "
        "which is why 'I wonder where your shoes are' does not fool it.")


def slide_biencoder(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "The Bi-Encoder: Fast First-Pass Matching")
    _bullets(s, [
        ("A model that reads one sentence at a time and turns it into a \"meaning fingerprint\" (a list of numbers).", 0, False, None),
        ("Both our library examples and each new sentence get a fingerprint; we compare them with a closeness score.", 0, False, None),
        ("\"Bi\" means two sentences are encoded separately, so the whole library is fingerprinted once, in advance.", 0, True, BLUE),
        ("Why it matters: it is extremely fast. A new sentence just gets compared against the stored fingerprints, so it scales to hours of transcript.", 0, True, GREEN),
        ("Trade-off: it judges each sentence on its own, so the hardest look-alikes can still slip through.", 0, True, AMBER),
        ("", 0, False, None),
        ("Like a librarian who has memorized the shelves and instantly points to the closest book.", 1, True, GREY),
    ], top=Inches(1.8), size=19, gap=11)
    _note(s,
        "This slide unpacks the first half of the detector. A bi-encoder is the matcher. "
        "It reads a single sentence and produces a 'meaning fingerprint', a list of numbers "
        "that summarizes what the sentence means. We do this once for every example in our "
        "trusted library, ahead of time, and then for each new sentence as it arrives. "
        "Classifying is then just measuring closeness between fingerprints. The key word is "
        "'bi', meaning the two sentences are encoded separately and never have to be "
        "processed together, which is exactly what makes it fast enough to run over hours "
        "of recordings. The honest trade-off, and the reason the next slide exists: because "
        "it scores each sentence on its own, the very hardest look-alikes can still fool it. "
        "That is what the re-ranker is there to catch.")


def slide_reranker(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "The Re-Ranker: A Careful Second Opinion")
    _bullets(s, [
        ("The bi-encoder is fast but shallow. The re-ranker is slow but careful.", 0, True, BLUE),
        ("It reads the new sentence and a candidate match together, so it can weigh how the words actually interact.", 0, False, None),
        ("We run it only on the borderline cases the bi-encoder flags, so we get the accuracy without paying the speed cost on every sentence.", 0, True, GREEN),
        ("Here the re-ranker is a language model acting as a gate: \"is this really science talk?\"", 0, True, AMBER),
        ("", 0, False, None),
        ("Net effect: the fast pass casts a wide net; the careful pass removes the false alarms. This is the two-stage design behind the Step 11 results.", 1, True, GREY),
    ], top=Inches(1.9), size=19, gap=12)
    _note(s,
        "This slide unpacks the second half of the detector and explains the 'second "
        "opinion' from the analogy slide. A re-ranker reads the new sentence and a "
        "candidate library match together, at the same time, so it can judge how the words "
        "interact rather than just comparing two separate fingerprints. That makes it more "
        "accurate, but also much slower, because nothing can be pre-computed. The trick is "
        "that we do not run it on everything. The fast bi-encoder casts a wide net and "
        "flags the uncertain cases, and only those go to the re-ranker for a careful call. "
        "In our system the re-ranker is a language model acting as a gate that answers, in "
        "plain language, 'is this really science talk?' Tie it back: this fast-then-careful "
        "pairing is exactly the 'fine-tuned bi-encoder plus LLM re-ranker' deployed system "
        "you will see in the Step 11 evaluation. Fast pass for recall, careful pass for "
        "precision.")


def slide_scale(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "Turning ~200 Examples into ~8,500 Training Rows")
    _table(s, [
        ["Source", "Count", "What it is"],
        ["Curated positives", "196", "Your expert-labeled science-talk examples"],
        ["Paraphrases", "516", "Same meaning, different speaking styles"],
        ["Non-science examples", "7,815", "Real classroom talk + tricky look-alikes"],
        ["Total training rows", "~8,500", "Stratified into train / val / test"],
    ], left=Inches(0.7), top=Inches(2.0), width=Inches(12.0), height=Inches(3.0),
       col_widths=[Inches(3.4), Inches(1.8), Inches(6.8)], font=16)
    _, tf = _textbox(s, Inches(0.7), Inches(5.4), Inches(12.0), Inches(1.2))
    r = tf.paragraphs[0].add_run()
    r.text = ("The \"tricky look-alikes\" (hard negatives) are the secret sauce. For "
              "example, \"I wonder where your shoes are\" looks like science but isn't.")
    _set_run(r, size=16, italic=True, color=AMBER)
    _note(s,
        "The challenge: we started with only ~200 expert-labeled examples, which is far "
        "too few to train a reliable detector. So we expanded responsibly. Paraphrases: "
        "we rephrase each real example into the different ways a teacher or parent might "
        "actually say it (circle time versus side-by-side play), same meaning, more "
        "phrasings. Non-science examples: we pulled real classroom talk and also generated "
        "'tricky look-alikes.' Stress the bottom line: every synthetic row is anchored to "
        "one of your real labeled examples. We are stretching your labels, not inventing "
        "new science out of thin air. The hard-negative example on screen usually gets a "
        "laugh and makes the point stick: surface words look scientific, meaning is not.")


def slide_qc(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "Built-in Quality Control")
    _bullets(s, [
        ("Every generated row gets a confidence score from three independent signals:", 0, True, NAVY),
        ("The AI's own self-rating", 1, False, None),
        ("A neutral meaning-similarity check", 1, False, None),
        ("Simple structural sanity rules", 1, False, None),
        ("", 0, False, None),
        ("Then each row is auto-sorted into one of three piles:", 0, True, NAVY),
        ("Auto: trusted, use directly", 1, True, GREEN),
        ("Spot: random spot-check", 1, True, AMBER),
        ("Review: needs human eyes", 1, True, BLUE),
    ], top=Inches(1.8), size=20, gap=8)
    _note(s,
        "Because we generate data, we cannot blindly trust it, so every generated row gets "
        "a confidence score built from three independent signals that have to agree: the "
        "AI's own self-rating, a neutral meaning-similarity check, and simple sanity rules. "
        "We combine them so that if any one signal is weak, the score drops, and we only "
        "auto-trust a row when all three agree. Then we sort every row into three piles: "
        "Auto (confident enough to use directly), Spot (use, but random spot-check), and "
        "Review (a human must look before we use it). One reassuring aside: the cut-offs "
        "live in a settings file, so when your audit told us to adjust, that was a "
        "config change, not a code rewrite. This slide sets up the audit results next.")


def slide_results(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "Your Audit: The Results",
               subtitle="50-row routing audit · 47 rows scored")
    _table(s, [
        ["Bucket", "Rows", "Agreement with you"],
        ["Auto (use directly)", "23", "100%"],
        ["Spot (spot-check)", "12", "58%"],
        ["Review (human eyes)", "12", "58%"],
        ["Overall", "47", "78.7%"],
    ], left=Inches(0.7), top=Inches(1.9), width=Inches(8.0), height=Inches(2.8),
       col_widths=[Inches(3.6), Inches(1.6), Inches(2.8)], font=16)

    box = s.shapes.add_shape(1, Inches(9.0), Inches(1.9), Inches(3.7), Inches(2.8))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.color.rgb = BLUE
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_right = Pt(10); tf.margin_top = Pt(10)
    r = tf.paragraphs[0].add_run()
    r.text = "Key finding"
    _set_run(r, size=16, bold=True, color=BLUE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = ("The system never wrongly auto-accepted anything. Every "
              "disagreement was in the spot/review piles, which already get "
              "human attention.")
    _set_run(r, size=14, color=NAVY)

    _, tf2 = _textbox(s, Inches(0.7), Inches(5.1), Inches(12.0), Inches(1.6))
    r = tf2.paragraphs[0].add_run()
    r.text = ("Overall 78.7% is just under our 80% target, but the misses are "
              "concentrated exactly where we already planned to look.")
    _set_run(r, size=16, italic=True, color=GREY)
    p = tf2.add_paragraph()
    r = p.add_run()
    r.text = "By kind: paraphrases 100% · non-science examples 70%."
    _set_run(r, size=14, color=GREY)
    _note(s,
        "This is your data, the 50-row sample you reviewed, 47 of which were scored. Be "
        "candid and lead with the honest headline: overall agreement was 78.7%, just under "
        "our 80% target. But the structure of the result is the good news, so do not bury "
        "it. Every row the system put in the 'Auto / use-directly' pile matched your "
        "judgment, 100%, zero exceptions. All of the disagreement was in the Spot and "
        "Review piles, which by design already go to a human. So the model and you "
        "disagreed only about how much scrutiny a borderline row needs, never about "
        "whether the auto-trusted rows were right. By kind: paraphrases matched you 100%, "
        "the non-science examples 70%. If asked: the misses are concentrated on hard "
        "negatives, which is exactly the corpus limitation we hit on the caveat slide.")


def slide_meaning(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "What This Means")
    _bullets(s, [
        ("The \"use-it-blindly\" pile is reliable, a 100% match with your judgments.", 0, True, GREEN),
        ("Disagreements are about how much review a borderline row needs, not about mislabeling science vs. non-science.", 0, False, None),
        ("Riskier rows are correctly routed to humans, never silently trusted.", 0, True, BLUE),
        ("", 0, False, None),
        ("This is the conservative behavior we want, because false positives cost researcher time.", 1, True, GREY),
    ], top=Inches(2.0), size=21, gap=14)
    _note(s,
        "Interpret the audit for them so they do not have to. The 'use-it-blindly' pile is "
        "trustworthy. That is the part we would run without supervision, and it matched you "
        "perfectly. The disagreements are not science-vs-not-science mistakes; they are "
        "judgment calls about how much review a borderline item deserves. And critically, "
        "the riskier items always err toward a human rather than getting silently accepted. "
        "Land the principle: these are 'safe errors.' Given that a false positive costs a "
        "researcher's time, a system that over-routes uncertain cases to people is exactly "
        "the conservative behavior we want. Pause here before the harder, honest slides.")


def slide_eval(prs):
    s = _blank(prs); _bg(s)
    _title_bar(
        s, "Step 11: Evaluation Results",
        subtitle="DoD (mechanics): PASS  ·  deployed = fine-tuned bi-encoder + LLM re-ranker @ threshold 0.9",
    )

    # PASS badge
    badge = s.shapes.add_shape(1, Inches(10.7), Inches(0.5), Inches(2.0), Inches(0.7))
    badge.fill.solid(); badge.fill.fore_color.rgb = GREEN
    badge.line.fill.background()
    btf = badge.text_frame; btf.word_wrap = True
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = "DoD: PASS"
    _set_run(br, size=18, bold=True, color=WHITE)

    # Ablation table (compact)
    _table(s, [
        ["System", "Thr", "Test F1", "Test Recall", "Hard-slice Recall"],
        ["Frozen cosine", "0.65", "0.952", "1.000", "1.000"],
        ["Fine-tuned cosine", "0.40", "1.000", "1.000", "1.000"],
        ["Fine-tuned + re-rank (deployed)", "0.90", "0.947", "0.900", "0.867"],
    ], left=Inches(0.7), top=Inches(1.75), width=Inches(12.0), height=Inches(2.0),
       col_widths=[Inches(4.4), Inches(1.3), Inches(2.0), Inches(2.2), Inches(2.1)],
       font=14)

    # Metrics callout box
    box = s.shapes.add_shape(1, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.color.rgb = BLUE
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_top = Pt(8)
    r = tf.paragraphs[0].add_run()
    r.text = "Deployed system on held-out test:  Precision 1.00  ·  Recall 0.90  ·  F1 0.95"
    _set_run(r, size=17, bold=True, color=NAVY)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Confusion:  TP 9  ·  FP 0  ·  FN 1  ·  TN 10   =   zero false positives."
    _set_run(r, size=15, color=GREEN)

    # Footnote that bridges into the caveat
    _, tf2 = _textbox(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.6))
    r = tf2.paragraphs[0].add_run()
    r.text = ("DoD is on mechanics: tuned on dev (not test), recall floor "
              "enforced, ablation produced, per-row predictions written.")
    _set_run(r, size=14, italic=True, color=GREY)
    p = tf2.add_paragraph()
    r = p.add_run()
    r.text = ("Y1 non-science talk is trivially separable, so F1 saturates. The "
              "hard-informal slice recall (0.867) is the number to watch; revisit with Y2 distractors.")
    _set_run(r, size=14, italic=True, color=AMBER)
    _note(s,
        "Step 11 is our formal evaluation, separate from your manual audit. This is the "
        "model graded against held-out answers, and the headline is that it passed its "
        "Definition of Done. The deployed system (the fine-tuned matcher plus the "
        "re-ranker second opinion) got Precision 1.00, Recall 0.90, F1 0.95 on the test "
        "set. In plain terms, it made zero false alarms and missed only one true science "
        "moment out of ten. That one miss was an off-task 'clean up the yogurt spill' "
        "comment, which is a defensible miss. Now be intellectually honest about the "
        "ablation table: the simpler cosine-only systems score a perfect 1.0, but that is "
        "NOT because they are better. It is because Year-1 non-science talk (classroom "
        "management) is so different from science talk that almost anything separates "
        "them. The scores are flattered by easy data. That is why I point to the "
        "'hard-informal slice', a deliberately difficult subset, where the deployed "
        "system gets 0.867 recall. That harder number is the realistic one, and it is the "
        "number that will move when we hit Year-2's messier talk. 'DoD on mechanics' means "
        "we proved the evaluation was done correctly (tuned on dev not test, recall floor "
        "enforced, ablation and per-row predictions produced), not that the absolute score "
        "is final. This is the perfect hand-off to the caveat.")


def slide_caveat(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "The Caveat: Being Honest")
    _bullets(s, [
        ("The whole system rests on a small, skewed seed corpus:", 0, True, NAVY),
        ("Only 196 labeled utterances: 191 science vs. just 5 non-science.", 1, False, None),
        ("Sub-types are lopsided: content (171) & observation (68) dominate;", 1, False, None),
        ("prediction (7), causal reasoning (8), evidence (4) are rare.", 1, True, AMBER),
        ("Most non-science training data is synthetic or transcript-mined, not expert-labeled.", 1, False, None),
        ("", 0, False, None),
        ("The model is strongest where we have data, and thinnest on the rare practices.", 0, True, BLUE),
    ], top=Inches(1.8), size=20, gap=10)
    _note(s,
        "This is the slide to deliver slowly and without spin. Credibility comes from "
        "naming the limits ourselves. Everything we have shown rests on a small, lopsided "
        "seed corpus: 196 labeled utterances, and of those only 5 are non-science, so the "
        "model has seen very few hand-labeled examples of what to reject. The sub-types are "
        "also skewed. 'Content' and 'observation' are well represented, but 'prediction', "
        "'causal reasoning', and 'evidence' have only a handful of examples each, so the "
        "model is weakest exactly on those richer practices. And most of our non-science "
        "training data is synthetic or transcript-mined rather than expert-labeled. Net: "
        "the system is genuinely strong where we have data and thin where we do not, and "
        "more labels, especially of the rare practices and of real non-science talk, is "
        "the single highest-leverage thing that would improve it. That tees up the "
        "decision on the next slide.")


def slide_decision(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "The Decision Point")

    a = s.shapes.add_shape(1, Inches(0.7), Inches(1.9), Inches(5.8), Inches(4.4))
    a.fill.solid(); a.fill.fore_color.rgb = LIGHT_BLUE
    a.line.color.rgb = BLUE
    tf = a.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(14); tf.margin_right = Pt(14); tf.margin_top = Pt(12)
    r = tf.paragraphs[0].add_run()
    r.text = "Option A: Proceed with Y2 data"
    _set_run(r, size=20, bold=True, color=BLUE)
    for t in [
        "Apply this pipeline to the Year-2 data now",
        "Fastest to results; accept current limits on rare sub-types",
        "Good if we want a working system to iterate on",
    ]:
        p = tf.add_paragraph(); p.space_before = Pt(8)
        rr = p.add_run(); rr.text = "•  " + t
        _set_run(rr, size=15, color=NAVY)

    b = s.shapes.add_shape(1, Inches(6.85), Inches(1.9), Inches(5.8), Inches(4.4))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xE0)
    b.line.color.rgb = AMBER
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(14); tf.margin_right = Pt(14); tf.margin_top = Pt(12)
    r = tf.paragraphs[0].add_run()
    r.text = "Option B: Repeat with more labeled data"
    _set_run(r, size=20, bold=True, color=AMBER)
    for t in [
        "Gather more expert labels (esp. rare practices + more non-science positives)",
        "Re-run the same pipeline for a higher-quality, more balanced corpus",
        "Slower, but lifts accuracy before we scale",
    ]:
        p = tf.add_paragraph(); p.space_before = Pt(8)
        rr = p.add_run(); rr.text = "•  " + t
        _set_run(rr, size=15, color=NAVY)

    _, tf3 = _textbox(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.7))
    r = tf3.paragraphs[0].add_run()
    r.text = "Not mutually exclusive long-term, but which do we prioritize now?"
    _set_run(r, size=15, italic=True, color=GREY)
    _note(s,
        "Present the two paths even-handedly; this is genuinely her call. Option A is to "
        "point the existing pipeline at the Year-2 data now. This is the fastest way to "
        "results and a working system we can iterate on, and Year-2's harder, messier talk "
        "is also where the evaluation numbers become truly informative. The trade-off is "
        "we would be accepting today's weakness on the rare sub-types. Option B is to "
        "invest first in more labeled data, especially the rare practices and more genuine "
        "non-science examples, then re-run the exact same pipeline for a higher-quality, "
        "more balanced corpus. Slower, but it lifts accuracy before we scale up. Make clear "
        "these are not mutually exclusive over the long run. The real question is "
        "sequencing: what do we do first? Then hand the decision to her on the next slide.")


def slide_ask(prs):
    s = _blank(prs); _bg(s)
    _title_bar(s, "Recommendation / Ask")
    _bullets(s, [
        ("The infrastructure is done, tested, and reproducible (free re-runs via caching).", 0, True, GREEN),
        ("The pipeline is ready to ingest more labels with zero code changes.", 0, True, BLUE),
        ("", 0, False, None),
        ("Decision needed from you: prioritize Y2 data now (A), or more labeled data first (B)?", 0, True, NAVY),
    ], top=Inches(2.2), size=22, gap=16)
    _note(s,
        "Close with confidence and brevity. Two reassurances first: the whole pipeline is "
        "built, tested, and reproducible (re-runs are essentially free because results "
        "are cached), and it is ready to ingest more labels with zero code changes, so "
        "whichever path we choose, there is no rebuild cost. Then make the single explicit "
        "ask and stop talking: 'So the one decision I need from you today is which to "
        "prioritize: pointing this at Year-2 data now, or investing in more labeled data "
        "first.' If you have a recommendation, state it in one sentence here; otherwise "
        "hand it to her and let the silence do the work.")


def build(output_path: Path = OUTPUT_PATH) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_detector(prs)
    slide_biencoder(prs)
    slide_reranker(prs)
    slide_scale(prs)
    slide_qc(prs)
    slide_results(prs)
    slide_meaning(prs)
    slide_eval(prs)
    slide_caveat(prs)
    slide_decision(prs)
    slide_ask(prs)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
